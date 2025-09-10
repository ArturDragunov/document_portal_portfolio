from pathlib import Path
from typing import Iterable, List, Optional
import tempfile
import os
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import camelot
import easyocr
from langchain.schema import Document
from logger import GLOBAL_LOGGER as log


class EmbeddedContentExtractor:
    """Information extractor using EasyOCR to grab text from embedded images."""

    def __init__(self, lang: str = "en"):
        # Initialize EasyOCR once for efficiency
        self.reader = easyocr.Reader([lang], gpu=False)

    # --------------------------
    # Helpers
    # --------------------------
    def _write_temp_file(self, data: bytes, suffix: str = ".png") -> str:
        tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
        try:
            tmp.write(data)
            return tmp.name
        finally:
            tmp.close()

    def _extract_text_from_image(self, image_path: str) -> str:
        """Run EasyOCR on an image and return extracted text."""
        try:
            results = self.reader.readtext(image_path, detail=0)
            return "\n".join(results).strip()
        except Exception as e:
            log.error(f"OCR failed for {image_path}: {e}")
            return ""

    def process_embedded_image(self, image_data: bytes, metadata: dict) -> Optional[Document]:
        tmp_path = self._write_temp_file(image_data)
        try:
            content = self._extract_text_from_image(tmp_path)
            if content:
                return Document(page_content=content, metadata=metadata)
        finally:
            os.unlink(tmp_path)
        return None

    # --------------------------
    # Extractors
    # --------------------------
    def extract_images_from_pdf(self, pdf_path: str) -> List[Document]:
        """Extract embedded images from PDF pages and OCR them."""
        docs = []
        try:
            pdf_doc = fitz.open(pdf_path)
            for page_num, page in enumerate(pdf_doc, start=1):
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                img_data = pix.tobytes("png")
                doc_obj = self.process_embedded_image(
                    img_data,
                    {"source": pdf_path, "page": page_num, "type": "pdf_embedded_content"},
                )
                if doc_obj:
                    docs.append(doc_obj)
            pdf_doc.close()
        except Exception as e:
            log.error(f"Failed to extract from PDF {pdf_path}: {e}")
        return docs
    
    def extract_tables_from_pdf(self, pdf_path: str) -> List[Document]:
        """Extract tables from a PDF using Camelot and return as Documents."""
        docs = []
        try:
            # Try lattice first (works best for bordered tables), fallback to stream
            tables = camelot.read_pdf(pdf_path, flavor="lattice", pages="all")
            if len(tables) == 0:
                tables = camelot.read_pdf(pdf_path, flavor="stream", pages="all")

            for i, table in enumerate(tables):
                df = table.df  # pandas DataFrame
                if not df.empty:
                    docs.append(
                        Document(
                            page_content=df.to_csv(index=False),
                            metadata={
                                "source": pdf_path,
                                "type": "pdf_table",
                                "table_index": i,
                                "parsing_report": table.parsing_report,
                            },
                        )
                    )
        except Exception as e:
            log.error(f"Table extraction failed for {pdf_path}: {e}")
        return docs
    
    def extract_images_from_docx(self, docx_path: str) -> List[Document]:
        """Extract embedded images from DOCX and OCR them."""
        docs = []
        try:
            doc = DocxDocument(docx_path)
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        doc_obj = self.process_embedded_image(
                            rel.target_part.blob,
                            {"source": docx_path, "type": "docx_embedded_image", "image_ref": rel.target_ref},
                        )
                        if doc_obj:
                            docs.append(doc_obj)
                    except Exception as e:
                        log.error(f"Failed to process image from DOCX: {e}")
        except Exception as e:
            log.error(f"Failed to extract from DOCX {docx_path}: {e}")
        return docs

    def extract_tables_from_docx(self, docx_path: str) -> List[Document]:
        docs = []
        try:
            doc = DocxDocument(docx_path)
            for table_idx, table in enumerate(doc.tables, start=1):
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                csv_content = "\n".join([",".join(r) for r in rows])
                docs.append(
                    Document(
                        page_content=csv_content,
                        metadata={
                            "source": docx_path,
                            "type": "docx_table",
                            "table_index": table_idx,
                        },
                    )
                )
        except Exception as e:
            log.error(f"Failed to extract tables from DOCX {docx_path}: {e}")
        return docs

    def extract_images_from_pptx(self, pptx_path: str) -> List[Document]:
        """Extract embedded images from PPTX and OCR them."""
        docs = []
        try:
            prs = Presentation(pptx_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture type
                        try:
                            doc_obj = self.process_embedded_image(
                                shape.image.blob,
                                {"source": pptx_path, "slide": slide_num, "type": "pptx_embedded_image"},
                            )
                            if doc_obj:
                                docs.append(doc_obj)
                        except Exception as e:
                            log.error(f"Failed to process image from slide {slide_num}: {e}")
        except Exception as e:
            log.error(f"Failed to extract from PPTX {pptx_path}: {e}")
        return docs
    
    def extract_tables_from_pptx(self, pptx_path: str) -> List[Document]:
        docs = []
        try:
            prs = Presentation(pptx_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                for shape_num, shape in enumerate(slide.shapes, start=1):
                    if shape.has_table:
                        table = shape.table
                        rows = []
                        for row in table.rows:
                            rows.append([cell.text.strip() for cell in row.cells])
                        csv_content = "\n".join([",".join(r) for r in rows])
                        docs.append(
                            Document(
                                page_content=csv_content,
                                metadata={
                                    "source": pptx_path,
                                    "slide": slide_num,
                                    "type": "pptx_table",
                                    "table_index": shape_num,
                                },
                            )
                        )
        except Exception as e:
            log.error(f"Failed to extract tables from PPTX {pptx_path}: {e}")
        return docs